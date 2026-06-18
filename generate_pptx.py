#!/usr/bin/env python3
"""
Generate PowerPoint presentation from MVP roadmap
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Output directory
OUTPUT_DIR = "/Users/dhiazfathra/Documents/GitHub/dhiazfathra/figma-mcp/docs/superpowers/specs"

# Colors
COLORS = {
    'primary': RGBColor(0x2F, 0x54, 0x96),      # Dark Blue
    'secondary': RGBColor(0x44, 0x72, 0xC4),    # Blue
    'accent': RGBColor(0xED, 0x7D, 0x31),       # Orange
    'success': RGBColor(0x70, 0xAD, 0x47),      # Green
    'danger': RGBColor(0xFF, 0x00, 0x00),        # Red
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
    'light_gray': RGBColor(0xF2, 0xF2, 0xF2),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
}

# Sprint data
SPRINTS = [
    # Backend sprints
    {"sprint": "B-1", "week": 1, "start": "2026-06-22", "duration": 7, "focus": "Foundation", "deliverables": "Project setup, Auth APIs, User model, DB schema", "team": "Backend"},
    {"sprint": "B-2", "week": 2, "start": "2026-06-29", "duration": 7, "focus": "Core APIs", "deliverables": "Brand, Booking, Promotion APIs", "team": "Backend"},
    {"sprint": "B-3", "week": 3, "start": "2026-07-06", "duration": 7, "focus": "Extended APIs", "deliverables": "Notification, Profile, Search APIs", "team": "Backend"},
    {"sprint": "B-4", "week": 4, "start": "2026-07-13", "duration": 7, "focus": "Domain APIs", "deliverables": "Treatment, Clinic, Skin Analysis APIs", "team": "Backend"},
    {"sprint": "B-5", "week": 5, "start": "2026-07-20", "duration": 7, "focus": "Phase 2 APIs", "deliverables": "Healthy Food, Shop APIs", "team": "Backend"},
    {"sprint": "B-6", "week": 6, "start": "2026-07-27", "duration": 7, "focus": "Phase 2 APIs", "deliverables": "Treatment Transaction, Skin Journey APIs", "team": "Backend"},
    {"sprint": "B-7", "week": 7, "start": "2026-08-03", "duration": 7, "focus": "CMS", "deliverables": "Dashboard, Content management", "team": "Backend"},
    {"sprint": "B-8", "week": 8, "start": "2026-08-10", "duration": 7, "focus": "Polish", "deliverables": "Documentation, Optimization, Security", "team": "Backend"},
    {"sprint": "B-BUF", "week": "9-10", "start": "2026-08-17", "duration": 14, "focus": "Buffer", "deliverables": "Integration testing, Bug fixes, Security audit", "team": "Backend"},
    
    # Mobile sprints
    {"sprint": "M-1", "week": 3, "start": "2026-07-06", "duration": 7, "focus": "Setup & Auth", "deliverables": "Flutter project, Login screen, OTP verification", "team": "Mobile"},
    {"sprint": "M-2", "week": 4, "start": "2026-07-13", "duration": 7, "focus": "Auth Complete", "deliverables": "Register, Forgot Password, Onboarding tutorial", "team": "Mobile"},
    {"sprint": "BUF-1", "week": 5, "start": "2026-07-20", "duration": 7, "focus": "Buffer", "deliverables": "Auth integration testing, Bug fixes", "team": "Mobile"},
    {"sprint": "M-3", "week": 6, "start": "2026-07-27", "duration": 7, "focus": "Homepage", "deliverables": "Non-login homepage, Logged-in homepage", "team": "Mobile"},
    {"sprint": "M-4", "week": 7, "start": "2026-08-03", "duration": 7, "focus": "Homepage Complete", "deliverables": "Search functionality, Brand carousel", "team": "Mobile"},
    {"sprint": "M-5", "week": 8, "start": "2026-08-10", "duration": 7, "focus": "Brands", "deliverables": "Brand listing, Brand detail", "team": "Mobile"},
    {"sprint": "M-6", "week": 9, "start": "2026-08-17", "duration": 7, "focus": "Brands Complete", "deliverables": "Clinic page, Treatment page", "team": "Mobile"},
    {"sprint": "M-7", "week": 10, "start": "2026-08-24", "duration": 7, "focus": "Booking (Part 1)", "deliverables": "Clinic selection, Treatment selection, Date picker", "team": "Mobile"},
    {"sprint": "M-8", "week": 11, "start": "2026-08-31", "duration": 7, "focus": "Booking (Part 2)", "deliverables": "Time picker, Booking confirmation, Payment", "team": "Mobile"},
    {"sprint": "M-9", "week": 12, "start": "2026-09-07", "duration": 7, "focus": "Booking Complete", "deliverables": "Booking status, Push notifications", "team": "Mobile"},
    {"sprint": "BUF-2", "week": 13, "start": "2026-09-14", "duration": 7, "focus": "Buffer", "deliverables": "Booking integration testing, Payment validation", "team": "Mobile"},
    {"sprint": "M-10", "week": 14, "start": "2026-09-21", "duration": 7, "focus": "Profile", "deliverables": "Profile page, Edit profile, Address management", "team": "Mobile"},
    {"sprint": "M-11", "week": 15, "start": "2026-09-28", "duration": 7, "focus": "Profile Complete", "deliverables": "Change phone, Change password, Face ID, Privacy", "team": "Mobile"},
    {"sprint": "M-12", "week": 16, "start": "2026-10-05", "duration": 7, "focus": "Promotions", "deliverables": "Promotion listing, Promotion detail", "team": "Mobile"},
    {"sprint": "M-13", "week": 17, "start": "2026-10-12", "duration": 7, "focus": "History & Skin", "deliverables": "Booking history, Skin analysis questionnaire", "team": "Mobile"},
    {"sprint": "M-14", "week": 18, "start": "2026-10-19", "duration": 7, "focus": "Phase 2 Start", "deliverables": "Healthy Food listing, Food detail, Food cart", "team": "Mobile"},
    {"sprint": "M-15", "week": 19, "start": "2026-10-26", "duration": 7, "focus": "Food Complete", "deliverables": "Food order, Order tracking", "team": "Mobile"},
    {"sprint": "M-16", "week": 20, "start": "2026-11-02", "duration": 7, "focus": "Shop", "deliverables": "Shop listing, Product detail, Shop cart", "team": "Mobile"},
    {"sprint": "M-17", "week": 21, "start": "2026-11-09", "duration": 7, "focus": "Shop Complete", "deliverables": "Shop order, Order tracking", "team": "Mobile"},
    {"sprint": "M-18", "week": 22, "start": "2026-11-16", "duration": 7, "focus": "Treatment & Journey", "deliverables": "Treatment transaction, Skin journey", "team": "Mobile"},
    {"sprint": "BUF-3", "week": 23, "start": "2026-11-23", "duration": 7, "focus": "Buffer", "deliverables": "Phase 2 integration testing, Cross-feature validation", "team": "Mobile"},
    {"sprint": "M-19", "week": 24, "start": "2026-11-30", "duration": 7, "focus": "Polish (Part 1)", "deliverables": "UI polish, Animation, Performance", "team": "Mobile"},
    {"sprint": "M-20", "week": 25, "start": "2026-12-07", "duration": 7, "focus": "Polish (Part 2)", "deliverables": "Bug fixes, Edge cases, Accessibility", "team": "Mobile"},
    {"sprint": "M-21", "week": 26, "start": "2026-12-14", "duration": 7, "focus": "QA & Testing", "deliverables": "Integration testing, E2E testing", "team": "Mobile"},
    {"sprint": "BUF-4", "week": 27, "start": "2026-12-21", "duration": 7, "focus": "Buffer", "deliverables": "Final QA, Regression testing, App store prep", "team": "Mobile"},
    {"sprint": "M-22", "week": 28, "start": "2026-12-28", "duration": 7, "focus": "Store Prep", "deliverables": "App store assets, Submission, Release", "team": "Mobile"},
]

# Milestones
MILESTONES = [
    {"name": "Backend APIs Ready", "date": "2026-07-20", "week": 4},
    {"name": "Auth Complete", "date": "2026-07-20", "week": 4},
    {"name": "Homepage Live", "date": "2026-08-10", "week": 7},
    {"name": "Booking + Payment Live", "date": "2026-09-14", "week": 12},
    {"name": "Profile + Promotions Live", "date": "2026-10-19", "week": 17},
    {"name": "Phase 2 Live", "date": "2026-11-23", "week": 22},
    {"name": "App Store Submission", "date": "2027-01-04", "week": 28},
]

# Features
FEATURES = [
    {"name": "Authentication & Onboarding", "epic": 1, "priority": "P0", "sprint": "M-1 to M-2", "weeks": "3-4"},
    {"name": "Homepage", "epic": 2, "priority": "P0", "sprint": "M-3 to M-4", "weeks": "6-7"},
    {"name": "Brand Pages", "epic": 3, "priority": "P0", "sprint": "M-5 to M-6", "weeks": "8-9"},
    {"name": "Booking Clinic", "epic": 4, "priority": "P0", "sprint": "M-7 to M-9", "weeks": "10-12"},
    {"name": "Profile Management", "epic": 5, "priority": "P1", "sprint": "M-10 to M-11", "weeks": "14-15"},
    {"name": "Notifications", "epic": 6, "priority": "P1", "sprint": "M-9", "weeks": "12"},
    {"name": "Promotions", "epic": 7, "priority": "P1", "sprint": "M-12", "weeks": "16"},
    {"name": "History", "epic": 8, "priority": "P1", "sprint": "M-13", "weeks": "17"},
    {"name": "Skin Analysis", "epic": 9, "priority": "P1", "sprint": "M-13", "weeks": "17"},
    {"name": "Healthy Food", "epic": 10, "priority": "P2", "sprint": "M-14 to M-15", "weeks": "18-19"},
    {"name": "Shop", "epic": 11, "priority": "P2", "sprint": "M-16 to M-17", "weeks": "20-21"},
    {"name": "Treatment Transaction", "epic": 12, "priority": "P2", "sprint": "M-18", "weeks": "22"},
    {"name": "Skin Journey", "epic": 13, "priority": "P2", "sprint": "M-18", "weeks": "22"},
]


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    top = Inches(4.2)
    height = Inches(1)
    
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, layout='bullets'):
    """Add a content slide with bullets or table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1)
    
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.2), width - Inches(1), height - Inches(0.4))
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Add content
    if layout == 'bullets':
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark_gray']
            p.space_after = Pt(12)
            p.level = 0
    
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1)
    
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.2), width - Inches(1), height - Inches(0.4))
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Calculate table dimensions
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.2)
    
    # Create table
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths if provided
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['secondary']
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS['white']
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Add rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = COLORS['dark_gray']
    
    return slide


def add_gantt_slide(prs, title, sprints, start_week, end_week):
    """Add a simplified Gantt chart slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1)
    
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.2), width - Inches(1), height - Inches(0.4))
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Calculate dimensions
    num_weeks = end_week - start_week + 1
    row_height = Inches(0.35)
    label_width = Inches(1.5)
    chart_left = Inches(2)
    chart_width = Inches(7.5)
    week_width = chart_width / num_weeks
    
    # Add week labels
    top = Inches(1.3)
    for week in range(num_weeks):
        left = chart_left + (week_width * week)
        box = slide.shapes.add_textbox(left, top, week_width, Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = f"W{start_week + week}"
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.CENTER
    
    # Add Gantt bars
    top = Inches(1.7)
    for sprint_idx, sprint in enumerate(sprints):
        row_top = top + (row_height * sprint_idx)
        
        # Add sprint label
        label_box = slide.shapes.add_textbox(Inches(0.2), row_top, label_width, row_height)
        p = label_box.text_frame.paragraphs[0]
        p.text = f"{sprint['sprint']}: {sprint['focus']}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['dark_gray']
        p.alignment = PP_ALIGN.LEFT
        label_box.text_frame.word_wrap = True
        
        # Calculate bar position
        sprint_start_week = sprint['week'] if isinstance(sprint['week'], int) else int(sprint['week'].split('-')[0])
        sprint_duration = sprint['duration'] // 7 if sprint['duration'] >= 7 else 1
        
        if sprint_start_week >= start_week and sprint_start_week <= end_week:
            bar_start = chart_left + (week_width * (sprint_start_week - start_week))
            bar_width = week_width * sprint_duration
            
            # Determine color
            if 'BUF' in sprint['sprint']:
                color = COLORS['success']
            elif sprint['team'] == 'Backend':
                color = COLORS['secondary']
            else:
                color = COLORS['accent']
            
            # Add bar
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_start, row_top + Inches(0.05), bar_width, row_height - Inches(0.1))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
    
    # Add legend
    legend_top = top + (row_height * len(sprints)) + Inches(0.3)
    legend_items = [
        ("Backend", COLORS['secondary']),
        ("Mobile", COLORS['accent']),
        ("Buffer", COLORS['success']),
    ]
    
    for i, (label, color) in enumerate(legend_items):
        left = Inches(2) + (Inches(2) * i)
        
        # Color box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, legend_top, Inches(0.3), Inches(0.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Label
        label_box = slide.shapes.add_textbox(left + Inches(0.4), legend_top - Inches(0.05), Inches(1.5), Inches(0.3))
        p = label_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['dark_gray']
    
    return slide


def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "Dermaesthetic Super Apps", 
        "MVP Development Roadmap\nJune 2026 - January 2027")
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "14 Features (11 Phase 1 + 3 Phase 2)",
        "28 Weeks Development (1 sprint = 1 week)",
        "4 Buffer Weeks (strategic placement)",
        "Team: 2 Flutter Devs + 2 Fullstack Devs",
        "Timeline: June 2026 → January 2027",
        "CTO Target: Beat by ~5 months",
        "Deferred: Laundry, Loyalty & Membership",
    ])
    
    # Slide 3: Team Structure
    add_table_slide(prs, "Team Structure", 
        ["Role", "Count", "Responsibility"],
        [
            ["Flutter Developer 1", "1", "Mobile app (Android + iOS)"],
            ["Flutter Developer 2", "1", "Mobile app (Android + iOS)"],
            ["Fullstack Developer 1", "1", "Backend APIs + CMS"],
            ["Fullstack Developer 2", "1", "Backend APIs + CMS"],
        ],
        [3, 1, 5])
    
    # Slide 4: Buffer Strategy
    add_table_slide(prs, "Buffer Strategy", 
        ["Buffer", "Week", "Purpose"],
        [
            ["BUF-1", "5", "Auth integration testing, Bug fixes"],
            ["BUF-2", "13", "Booking integration, Payment validation"],
            ["BUF-3", "23", "Phase 2 integration, Cross-feature validation"],
            ["BUF-4", "27", "Final QA, Regression testing, Store prep"],
            ["B-BUF", "9-10", "Backend integration, Security audit"],
        ],
        [1.5, 1.5, 6])
    
    # Slide 5: Features Overview
    add_table_slide(prs, "Features Overview (14 Features)", 
        ["Feature", "Priority", "Sprint", "Weeks"],
        [[f["name"], f["priority"], f["sprint"], f["weeks"]] for f in FEATURES],
        [3, 1, 2.5, 2.5])
    
    # Slide 6: Backend Sprints
    backend_sprints = [s for s in SPRINTS if s['team'] == 'Backend']
    add_table_slide(prs, "Backend Sprints", 
        ["Sprint", "Week", "Focus", "Deliverables"],
        [[s["sprint"], str(s["week"]), s["focus"], s["deliverables"]] for s in backend_sprints],
        [1.5, 1, 2, 4.5])
    
    # Slide 7: Mobile Sprints (Part 1)
    mobile_sprints_1 = [s for s in SPRINTS if s['team'] == 'Mobile'][:12]
    add_table_slide(prs, "Mobile Sprints (Part 1)", 
        ["Sprint", "Week", "Focus", "Deliverables"],
        [[s["sprint"], str(s["week"]), s["focus"], s["deliverables"]] for s in mobile_sprints_1],
        [1.5, 1, 2, 4.5])
    
    # Slide 8: Mobile Sprints (Part 2)
    mobile_sprints_2 = [s for s in SPRINTS if s['team'] == 'Mobile'][12:]
    add_table_slide(prs, "Mobile Sprints (Part 2)", 
        ["Sprint", "Week", "Focus", "Deliverables"],
        [[s["sprint"], str(s["week"]), s["focus"], s["deliverables"]] for s in mobile_sprints_2],
        [1.5, 1, 2, 4.5])
    
    # Slide 9: Gantt Chart (Weeks 1-14)
    add_gantt_slide(prs, "Gantt Chart (Weeks 1-14)", 
        [s for s in SPRINTS if isinstance(s['week'], int) and s['week'] <= 14],
        1, 14)
    
    # Slide 10: Gantt Chart (Weeks 15-28)
    add_gantt_slide(prs, "Gantt Chart (Weeks 15-28)", 
        [s for s in SPRINTS if isinstance(s['week'], int) and s['week'] >= 15],
        15, 28)
    
    # Slide 11: Milestones
    add_table_slide(prs, "Milestones", 
        ["Milestone", "Date", "Week"],
        [[m["name"], m["date"], str(m["week"])] for m in MILESTONES],
        [4, 3, 2])
    
    # Slide 12: Risk Register
    add_table_slide(prs, "Risk Register", 
        ["Risk", "Impact", "Mitigation"],
        [
            ["Flutter dev bottleneck", "High", "Backend APIs ready 1-2 weeks ahead"],
            ["Doku payment issues", "Medium", "Early integration testing"],
            ["App store rejection", "Medium", "Follow guidelines, Pre-submission review"],
            ["Scope creep", "High", "Strict MVP scope, Deferred features documented"],
            ["Team burnout", "Medium", "Realistic sprint loads, Buffer weeks"],
        ],
        [3, 1.5, 4.5])
    
    # Slide 13: Next Steps
    add_content_slide(prs, "Next Steps", [
        "1. Review and approve this roadmap",
        "2. Set up Flutter project structure",
        "3. Set up backend API architecture",
        "4. Begin Sprint B-1: Foundation (June 22)",
        "5. Begin Sprint M-1: Setup & Auth (July 6)",
        "6. Weekly sprint reviews and adjustments",
    ])
    
    # Slide 14: Thank You
    add_title_slide(prs, 
        "Thank You", 
        "Questions?\nContact: Project Team")
    
    # Save presentation
    pptx_path = os.path.join(OUTPUT_DIR, 'dermaesthetic-mvp-roadmap.pptx')
    prs.save(pptx_path)
    print(f'PowerPoint saved: {pptx_path}')
    return pptx_path


if __name__ == '__main__':
    print('Generating PowerPoint presentation...')
    pptx_path = create_presentation()
    print(f'Done! File: {pptx_path}')
